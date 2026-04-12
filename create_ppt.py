from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1B, 0x22)
ACCENT_BLUE = RGBColor(0x58, 0xA6, 0xFF)
ACCENT_PURPLE = RGBColor(0xBC, 0x8C, 0xFF)
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_ORANGE = RGBColor(0xFF, 0x7B, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
LIGHT_GRAY = RGBColor(0xC9, 0xD1, 0xD9)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=WHITE, icon_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        # Icon
        run = p.add_run()
        run.text = "▸ "
        run.font.size = Pt(font_size)
        run.font.color.rgb = icon_color
        run.font.name = "Segoe UI"
        run.font.bold = True
        # Text
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = "Segoe UI"
    return txBox

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide1, BG_DARK)

# Decorative top bar
bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

# Emoji icons
add_text(slide1, Inches(4.5), Inches(1.5), Inches(4.5), Inches(1), "🐳  🔧  ☁️  🐙", 
         font_size=48, color=WHITE, alignment=PP_ALIGN.CENTER)

# Title
add_text(slide1, Inches(1.5), Inches(2.5), Inches(10.5), Inches(1.2), "TRIONIX", 
         font_size=60, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide1, Inches(1.5), Inches(3.5), Inches(10.5), Inches(0.8), 
         "CI/CD Pipeline: Containerization & Cloud Deployment", 
         font_size=28, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Divider line
divider = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.5), Inches(3.5), Pt(2))
divider.fill.solid()
divider.fill.fore_color.rgb = ACCENT_PURPLE
divider.line.fill.background()

# Tech stack
add_text(slide1, Inches(1.5), Inches(5), Inches(10.5), Inches(0.6), 
         "Docker  •  Jenkins  •  Azure Container Registry  •  Azure App Service  •  GitHub", 
         font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# Names
add_text(slide1, Inches(1.5), Inches(6), Inches(10.5), Inches(0.6),
         "Aditya Rana  |  Krishna Menon",
         font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: What is CI/CD? Problem Statement
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)

add_text(slide2, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Why CI/CD?", 
         font_size=36, color=ACCENT_BLUE, bold=True)
add_text(slide2, Inches(0.8), Inches(1.0), Inches(6), Inches(0.5), "The Problem & Our Solution", 
         font_size=18, color=GRAY)

# Problem Card
add_shape(slide2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), BG_CARD, RGBColor(0xFF, 0x45, 0x45))
add_text(slide2, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5), "❌  The Problem (Manual Deployment)", 
         font_size=20, color=RGBColor(0xFF, 0x45, 0x45), bold=True)
add_bullet_text(slide2, Inches(1.1), Inches(2.7), Inches(5), Inches(3.5), [
    "Manually build Docker image every time",
    "Manually log into Azure and push the image",
    "Manually restart the web server",
    "Prone to human errors and inconsistency",
    "Time-consuming: 15+ minutes per deploy"
], font_size=15, color=LIGHT_GRAY, icon_color=RGBColor(0xFF, 0x45, 0x45))

# Solution Card
add_shape(slide2, Inches(7), Inches(1.8), Inches(5.5), Inches(4.8), BG_CARD, ACCENT_GREEN)
add_text(slide2, Inches(7.3), Inches(2.0), Inches(5), Inches(0.5), "✅  Our Solution (Automated CI/CD)", 
         font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_text(slide2, Inches(7.3), Inches(2.7), Inches(5), Inches(3.5), [
    "One-click deployment via Jenkins",
    "Docker auto-packages the entire ML app",
    "Jenkins auto-pushes to Azure Registry",
    "Webhook auto-restarts the live website",
    "Deployment time: ~20 seconds + auto"
], font_size=15, color=LIGHT_GRAY, icon_color=ACCENT_GREEN)

# ============================================================
# SLIDE 3: Technology Stack
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)

add_text(slide3, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Technology Stack", 
         font_size=36, color=ACCENT_BLUE, bold=True)
add_text(slide3, Inches(0.8), Inches(1.0), Inches(8), Inches(0.5), "Four technologies working together", 
         font_size=18, color=GRAY)

# Card 1 - Docker
add_shape(slide3, Inches(0.5), Inches(1.8), Inches(2.8), Inches(4.8), BG_CARD, ACCENT_BLUE)
add_text(slide3, Inches(0.7), Inches(2.0), Inches(2.5), Inches(0.5), "🐳 Docker", 
         font_size=22, color=ACCENT_BLUE, bold=True)
add_text(slide3, Inches(0.7), Inches(2.6), Inches(2.5), Inches(0.4), "CONTAINERIZATION", 
         font_size=11, color=GRAY, bold=True)
add_bullet_text(slide3, Inches(0.7), Inches(3.2), Inches(2.4), Inches(3), [
    "Packages entire ML app",
    "Python + PyTorch + FFmpeg",
    "9.88 GB production image",
    "Runs identically everywhere"
], font_size=13, color=LIGHT_GRAY, icon_color=ACCENT_BLUE)

# Card 2 - Jenkins
add_shape(slide3, Inches(3.7), Inches(1.8), Inches(2.8), Inches(4.8), BG_CARD, ACCENT_ORANGE)
add_text(slide3, Inches(3.9), Inches(2.0), Inches(2.5), Inches(0.5), "🔧 Jenkins", 
         font_size=22, color=ACCENT_ORANGE, bold=True)
add_text(slide3, Inches(3.9), Inches(2.6), Inches(2.5), Inches(0.4), "CI/CD AUTOMATION", 
         font_size=11, color=GRAY, bold=True)
add_bullet_text(slide3, Inches(3.9), Inches(3.2), Inches(2.4), Inches(3), [
    "Runs on local machine",
    "3-stage pipeline script",
    "Checkout → Build → Push",
    "One-click deployment"
], font_size=13, color=LIGHT_GRAY, icon_color=ACCENT_ORANGE)

# Card 3 - Azure
add_shape(slide3, Inches(6.9), Inches(1.8), Inches(2.8), Inches(4.8), BG_CARD, ACCENT_PURPLE)
add_text(slide3, Inches(7.1), Inches(2.0), Inches(2.5), Inches(0.5), "☁️ Azure", 
         font_size=22, color=ACCENT_PURPLE, bold=True)
add_text(slide3, Inches(7.1), Inches(2.6), Inches(2.5), Inches(0.4), "CLOUD HOSTING", 
         font_size=11, color=GRAY, bold=True)
add_bullet_text(slide3, Inches(7.1), Inches(3.2), Inches(2.4), Inches(3), [
    "Container Registry (ACR)",
    "App Service (B1 tier)",
    "Auto-deploy webhooks",
    "Public HTTPS URL"
], font_size=13, color=LIGHT_GRAY, icon_color=ACCENT_PURPLE)

# Card 4 - GitHub
add_shape(slide3, Inches(10.1), Inches(1.8), Inches(2.8), Inches(4.8), BG_CARD, ACCENT_GREEN)
add_text(slide3, Inches(10.3), Inches(2.0), Inches(2.5), Inches(0.5), "🐙 GitHub", 
         font_size=22, color=ACCENT_GREEN, bold=True)
add_text(slide3, Inches(10.3), Inches(2.6), Inches(2.5), Inches(0.4), "SOURCE CONTROL", 
         font_size=11, color=GRAY, bold=True)
add_bullet_text(slide3, Inches(10.3), Inches(3.2), Inches(2.4), Inches(3), [
    "Central code repository",
    "Version control (Git)",
    "Stores CI/CD configs",
    "Pull Request workflow"
], font_size=13, color=LIGHT_GRAY, icon_color=ACCENT_GREEN)

# ============================================================
# SLIDE 4: Pipeline Flow
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_DARK)

add_text(slide4, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Pipeline Flow", 
         font_size=36, color=ACCENT_BLUE, bold=True)
add_text(slide4, Inches(0.8), Inches(1.0), Inches(8), Inches(0.5), "From code change to live website in one click", 
         font_size=18, color=GRAY)

# Flow boxes
flow_items = [
    ("1", "📝 Code Change", "Developer updates\ncode on laptop", ACCENT_BLUE),
    ("2", "🐙 Git Push", "Code pushed to\nGitHub repository", ACCENT_GREEN),
    ("3", "🔧 Jenkins Build", "Docker image built\nautomatically", ACCENT_ORANGE),
    ("4", "📦 Push to ACR", "Image uploaded to\nAzure Registry", ACCENT_PURPLE),
    ("5", "🔔 Webhook", "Registry notifies\nApp Service", RGBColor(0xFF, 0xD7, 0x00)),
    ("6", "🌐 Live!", "Website updates\nautomatically", ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(flow_items):
    x = Inches(0.5 + i * 2.1)
    y = Inches(2.2)
    
    # Card
    add_shape(slide4, x, y, Inches(1.9), Inches(3), BG_CARD, color)
    
    # Step number circle
    circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), y + Inches(0.2), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = BG_DARK
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Segoe UI"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False
    
    # Title
    add_text(slide4, x + Inches(0.1), y + Inches(0.9), Inches(1.7), Inches(0.5), title, 
             font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Description
    add_text(slide4, x + Inches(0.1), y + Inches(1.5), Inches(1.7), Inches(1.2), desc, 
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    
    # Arrow between boxes
    if i < len(flow_items) - 1:
        add_text(slide4, x + Inches(1.85), y + Inches(1.1), Inches(0.3), Inches(0.5), "→", 
                 font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)

# Bottom highlight
add_shape(slide4, Inches(2), Inches(5.8), Inches(9.5), Inches(1), BG_CARD, ACCENT_BLUE)
add_text(slide4, Inches(2.3), Inches(5.95), Inches(9), Inches(0.7), 
         "⚡ Total deployment time: ~20 seconds build  +  ~10 minutes Azure cold start  =  Fully Automated!", 
         font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: Key Files Created
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_DARK)

add_text(slide5, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "Key Files & Configuration", 
         font_size=36, color=ACCENT_BLUE, bold=True)

# Dockerfile card
add_shape(slide5, Inches(0.5), Inches(1.4), Inches(6), Inches(2.5), BG_CARD, ACCENT_BLUE)
add_text(slide5, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4), "📄 Dockerfile", 
         font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_text(slide5, Inches(0.8), Inches(2.2), Inches(5.5), Inches(1.5), [
    "Base: Python 3.10-slim",
    "System deps: FFmpeg, libzbar0",
    "Installs all Python packages from requirements.txt",
    "Runs Django server on port 8000"
], font_size=14, color=LIGHT_GRAY, icon_color=ACCENT_BLUE)

# Jenkinsfile card
add_shape(slide5, Inches(6.8), Inches(1.4), Inches(6), Inches(2.5), BG_CARD, ACCENT_ORANGE)
add_text(slide5, Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.4), "📄 Jenkinsfile", 
         font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_text(slide5, Inches(7.1), Inches(2.2), Inches(5.5), Inches(1.5), [
    "Declarative pipeline (3 stages)",
    "Stage 1: Checkout code from GitHub",
    "Stage 2: Build Docker image locally",
    "Stage 3: Push image to Azure ACR"
], font_size=14, color=LIGHT_GRAY, icon_color=ACCENT_ORANGE)

# Azure Config card
add_shape(slide5, Inches(0.5), Inches(4.3), Inches(6), Inches(2.5), BG_CARD, ACCENT_PURPLE)
add_text(slide5, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.4), "☁️ Azure Configuration", 
         font_size=20, color=ACCENT_PURPLE, bold=True)
add_bullet_text(slide5, Inches(0.8), Inches(5.1), Inches(5.5), Inches(1.5), [
    "Container Registry: trionix.azurecr.io (Basic tier)",
    "App Service: B1 tier, Linux, Central India",
    "Continuous Deployment via Webhooks",
    "WEBSITES_PORT=8000, START_TIME_LIMIT=1800"
], font_size=14, color=LIGHT_GRAY, icon_color=ACCENT_PURPLE)

# Design Decision card
add_shape(slide5, Inches(6.8), Inches(4.3), Inches(6), Inches(2.5), BG_CARD, ACCENT_GREEN)
add_text(slide5, Inches(7.1), Inches(4.5), Inches(5.5), Inches(0.4), "💡 Key Design Decisions", 
         font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_text(slide5, Inches(7.1), Inches(5.1), Inches(5.5), Inches(1.5), [
    "Local Jenkins → saves Azure credits",
    "B1 tier → enough RAM for PyTorch",
    "Webhooks → fully automated deployment",
    "os.path.join → cross-platform paths"
], font_size=14, color=LIGHT_GRAY, icon_color=ACCENT_GREEN)

# ============================================================
# SLIDE 6: Live Demo / Thank You
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, BG_DARK)

# Decorative bar
bar = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_PURPLE
bar.line.fill.background()

add_text(slide6, Inches(1.5), Inches(1.5), Inches(10.5), Inches(1), "Live Demo & Results", 
         font_size=48, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Results card
add_shape(slide6, Inches(2.5), Inches(3), Inches(8.5), Inches(3), BG_CARD, ACCENT_BLUE)

add_bullet_text(slide6, Inches(3), Inches(3.3), Inches(7.5), Inches(2.5), [
    "✅  Django ML app containerized with Docker (9.88 GB image)",
    "✅  Automated CI/CD pipeline with Jenkins (4 successful builds)",
    "✅  Deployed to Azure App Service (Central India region)",
    "✅  Auto-deployment via Container Registry webhooks",
    "✅  Live URL: trionix-app-xxxxx.azurewebsites.net"
], font_size=17, color=WHITE, icon_color=ACCENT_GREEN)

add_text(slide6, Inches(1.5), Inches(6.3), Inches(10.5), Inches(0.6), "Thank You! 🎉", 
         font_size=32, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Trionix_CICD_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
